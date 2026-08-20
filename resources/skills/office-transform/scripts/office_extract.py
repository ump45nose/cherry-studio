#!/usr/bin/env python3
"""Extract an anchored region of an Office/PDF file into a NEW file.

The source file is opened read-only and never modified. Anchors address the
document's own structural coordinates (worksheet range, body-level paragraph
ordinal, page number) — see SKILL.md for the anchor JSON shapes.

Format-specific third-party readers are imported lazily, so run this with the
dependency matching the source format, e.g.:

    uv run --with openpyxl python office_extract.py \
        --file /abs/report.xlsx \
        --anchor '{"format":"xlsx","sheet":"Sheet1","range":"A1:C10"}' \
        --out /abs/report-extract.csv

Dependencies by source format: xlsx -> openpyxl, docx -> python-docx,
pdf -> pypdf, pptx -> python-pptx.
"""

import argparse
import contextlib
import csv
import json
import os
import re
import sys
import tempfile
import zipfile
from pathlib import Path

A1_CELL_RE = re.compile(r"^([A-Z]{1,3})([1-9][0-9]*)$")

MAX_RANGE_CELLS = 1_000_000

# The SpreadsheetML grid (ECMA-376): columns A..XFD, rows 1..1048576.
MAX_COLUMN_INDEX = 16_384
MAX_ROW_NUMBER = 1_048_576
MAX_ZIP_ENTRIES = 10_000
MAX_ENTRY_BYTES = 256 * 1024 * 1024
MAX_TOTAL_BYTES = 1024 * 1024 * 1024


def fail(message: str) -> "sys.NoReturn":
    print(f"error: {message}", file=sys.stderr)
    raise SystemExit(1)


@contextlib.contextmanager
def atomic_output(out_path: Path):
    """Yield a staging path renamed onto `out_path` only on success, so a failure leaves nothing behind.

    Without this, an interrupted write leaves a partial file that both looks like a result and blocks
    the retry with "output path already exists".
    """
    handle, staging_name = tempfile.mkstemp(dir=out_path.parent, prefix=f".{out_path.name}.", suffix=".part")
    os.close(handle)
    staging = Path(staging_name)
    try:
        yield staging
        staging.replace(out_path)
    except BaseException:
        staging.unlink(missing_ok=True)
        raise


def column_to_index(letters: str) -> int:
    index = 0
    for char in letters:
        index = index * 26 + (ord(char) - ord("A") + 1)
    return index


def parse_a1_cell(ref: str) -> tuple[int, int]:
    match = A1_CELL_RE.match(ref)
    if not match:
        fail(f"invalid A1 cell reference: {ref!r}")
    column, row = column_to_index(match.group(1)), int(match.group(2))
    if column > MAX_COLUMN_INDEX or row > MAX_ROW_NUMBER:
        fail(f"cell {ref!r} is outside the worksheet grid (max XFD{MAX_ROW_NUMBER})")
    return column, row


def parse_a1_range(ref: str) -> tuple[int, int, int, int]:
    """Return (min_col, min_row, max_col, max_row) from 'B2' or 'A1:C10'."""
    parts = ref.split(":")
    if len(parts) > 2:
        fail(f"invalid A1 range: {ref!r}")
    start = parse_a1_cell(parts[0])
    end = parse_a1_cell(parts[-1])
    return (
        min(start[0], end[0]),
        min(start[1], end[1]),
        max(start[0], end[0]),
        max(start[1], end[1]),
    )


def preflight_zip(path: "Path") -> None:
    """Refuse pathological OOXML packages before a reader decompresses them."""
    try:
        with zipfile.ZipFile(path) as archive:
            infos = archive.infolist()
    except zipfile.BadZipFile:
        fail(f"not a valid OOXML package: {path}")
    if len(infos) > MAX_ZIP_ENTRIES:
        fail(f"package has {len(infos)} entries (limit {MAX_ZIP_ENTRIES})")
    total = 0
    for info in infos:
        if info.file_size > MAX_ENTRY_BYTES:
            fail(f"package entry {info.filename!r} decompresses to {info.file_size} bytes (limit {MAX_ENTRY_BYTES})")
        total += info.file_size
    if total > MAX_TOTAL_BYTES:
        fail(f"package decompresses to {total} bytes in total (limit {MAX_TOTAL_BYTES})")


def require_index(value, field: str, minimum: int) -> int:
    """Return an anchor ordinal, refusing anything int() would silently reinterpret.

    Bare int() accepts "3", 3.7 (truncated) and True (1), each of which addresses a different
    paragraph/page/slide than the caller meant and reports nothing. bool is checked first because it
    is a subclass of int.
    """
    if isinstance(value, bool) or not isinstance(value, int):
        fail(f"{field} must be an integer, not {type(value).__name__}: {value!r}")
    if value < minimum:
        fail(f"{field} must be >= {minimum}: {value!r}")
    return value


def slice_char_range(text: str, char_range) -> str:
    if char_range is None:
        return text
    # int() would happily take "26" (as two characters) or 1.9 (truncated), each of which slices a
    # different span than the caller asked for and reports nothing.
    if not isinstance(char_range, (list, tuple)) or len(char_range) != 2:
        fail(f"charRange must be a two-element [start, end] array: {char_range!r}")
    if not all(isinstance(bound, int) and not isinstance(bound, bool) for bound in char_range):
        fail(f"charRange bounds must be integers: {char_range!r}")
    start, end = char_range
    if start > end or start < 0:
        fail(f"invalid charRange: {char_range!r}")
    return text[start:end]


def cell_display(value) -> str:
    if value is None:
        return ""
    return str(value)


def write_markdown_table(rows: list[list[str]], out_path: Path) -> None:
    if not rows:
        fail("selection produced no rows")
    width = max(len(row) for row in rows)
    normalized = [row + [""] * (width - len(row)) for row in rows]
    escaped = [[cell.replace("|", "\\|").replace("\n", " ") for cell in row] for row in normalized]
    lines = ["| " + " | ".join(escaped[0]) + " |", "| " + " | ".join(["---"] * width) + " |"]
    lines.extend("| " + " | ".join(row) + " |" for row in escaped[1:])
    out_path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def extract_xlsx(src: Path, anchor: dict, out_path: Path, out_format: str) -> None:
    try:
        from openpyxl import Workbook, load_workbook
    except ImportError:
        fail("openpyxl is required for xlsx sources — rerun via `uv run --with openpyxl python ...`")

    sheet_name = anchor.get("sheet")
    range_ref = anchor.get("range")
    if not sheet_name or not range_ref:
        fail("xlsx anchor requires 'sheet' and 'range'")

    min_col, min_row, max_col, max_row = parse_a1_range(range_ref)
    area = (max_row - min_row + 1) * (max_col - min_col + 1)
    if area > MAX_RANGE_CELLS:
        fail(f"range {range_ref!r} covers {area} cells (limit {MAX_RANGE_CELLS}); select a smaller region")

    workbook = load_workbook(src, data_only=True, read_only=True)
    if sheet_name not in workbook.sheetnames:
        fail(f"worksheet not found: {sheet_name!r} (has: {workbook.sheetnames})")
    worksheet = workbook[sheet_name]

    values = [
        [cell.value for cell in row]
        for row in worksheet.iter_rows(min_row=min_row, max_row=max_row, min_col=min_col, max_col=max_col)
    ]
    workbook.close()

    if out_format == "xlsx":
        derived = Workbook()
        derived_sheet = derived.active
        derived_sheet.title = sheet_name[:31]
        for row in values:
            derived_sheet.append(row)
        derived.save(out_path)
    elif out_format == "csv":
        with out_path.open("w", newline="", encoding="utf-8") as handle:
            writer = csv.writer(handle)
            writer.writerows([[cell_display(value) for value in row] for row in values])
    elif out_format == "md":
        write_markdown_table([[cell_display(value) for value in row] for row in values], out_path)
    else:
        fail(f"unsupported output format for xlsx source: {out_format!r} (use xlsx, csv, or md)")


def extract_docx(src: Path, anchor: dict, out_path: Path, out_format: str) -> None:
    try:
        import docx
    except ImportError:
        fail("python-docx is required for docx sources — rerun via `uv run --with python-docx python ...`")

    if anchor.get("paragraph") is None:
        fail("docx anchor requires a non-negative 'paragraph' ordinal")
    paragraph_index = require_index(anchor.get("paragraph"), "docx anchor 'paragraph'", 0)

    document = docx.Document(str(src))
    paragraphs = document.paragraphs
    if paragraph_index >= len(paragraphs):
        fail(f"paragraph {paragraph_index} out of range (document has {len(paragraphs)} body paragraphs)")

    para_id = anchor.get("paraId")
    if para_id:
        def p_para_id(paragraph):
            for key, value in paragraph._p.attrib.items():
                if key.rsplit("}", 1)[-1] == "paraId":
                    return value
            return None

        matches = [i for i, p in enumerate(paragraphs) if p_para_id(p) == para_id]
        if len(matches) > 1:
            fail(f"paraId {para_id!r} matches {len(matches)} paragraphs; refusing an ambiguous anchor")
        if matches and matches[0] != paragraph_index:
            fail(
                f"paraId {para_id!r} resolves to paragraph {matches[0]} but the anchor says {paragraph_index} — "
                "the document changed since the anchor was captured; re-select"
            )
    text = slice_char_range(paragraphs[paragraph_index].text, anchor.get("charRange"))

    if out_format in ("txt", "md"):
        out_path.write_text(text + "\n", encoding="utf-8")
    elif out_format == "docx":
        derived = docx.Document()
        derived.add_paragraph(text)
        derived.save(str(out_path))
    else:
        fail(f"unsupported output format for docx source: {out_format!r} (use txt, md, or docx)")


def extract_pdf(src: Path, anchor: dict, out_path: Path, out_format: str) -> None:
    try:
        from pypdf import PdfReader, PdfWriter
    except ImportError:
        fail("pypdf is required for pdf sources — rerun via `uv run --with pypdf python ...`")

    if anchor.get("page") is None:
        fail("pdf anchor requires a one-based 'page' number")
    page_number = require_index(anchor.get("page"), "pdf anchor 'page'", 1)
    page_index = page_number - 1

    # OOXML sources go through preflight_zip; PDFs had no ceiling at all before PdfReader parsed them.
    source_bytes = src.stat().st_size
    if source_bytes > MAX_ENTRY_BYTES:
        fail(f"pdf is {source_bytes} bytes (limit {MAX_ENTRY_BYTES}); ask for a smaller file")

    reader = PdfReader(str(src))
    if page_index >= len(reader.pages):
        fail(f"page {page_number} out of range (document has {len(reader.pages)} pages)")
    page = reader.pages[page_index]

    if out_format == "pdf":
        writer = PdfWriter()
        writer.add_page(page)
        with out_path.open("wb") as handle:
            writer.write(handle)
    elif out_format in ("txt", "md"):
        text = slice_char_range(page.extract_text() or "", anchor.get("charRange"))
        out_path.write_text(text + "\n", encoding="utf-8")
    else:
        fail(f"unsupported output format for pdf source: {out_format!r} (use pdf, txt, or md)")


def iter_shapes_recursive(shapes):
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes_recursive(shape.shapes)


def shape_text_lines(shape) -> list[str]:
    if shape.has_text_frame:
        return [paragraph.text for paragraph in shape.text_frame.paragraphs]
    if getattr(shape, "has_table", False) and shape.has_table:
        return [" | ".join(cell.text for cell in row.cells) for row in shape.table.rows]
    return []


def extract_pptx(src: Path, anchor: dict, out_path: Path, out_format: str) -> None:
    try:
        from pptx import Presentation
    except ImportError:
        fail("python-pptx is required for pptx sources — rerun via `uv run --with python-pptx python ...`")

    if anchor.get("slide") is None:
        fail("pptx anchor requires a one-based 'slide' number")
    slide_number = require_index(anchor.get("slide"), "pptx anchor 'slide'", 1)
    slide_index = slide_number - 1

    presentation = Presentation(str(src))
    slides = list(presentation.slides)
    if slide_index >= len(slides):
        fail(f"slide {slide_number} out of range (deck has {len(slides)} slides)")
    slide = slides[slide_index]

    node_id = anchor.get("nodeId")
    if node_id is None:
        if anchor.get("paragraph") is not None or anchor.get("tableCell") is not None:
            fail("pptx anchor has 'paragraph'/'tableCell' but no 'nodeId'; refusing to fall back to whole-slide extraction")
        lines = [line for shape in iter_shapes_recursive(slide.shapes) for line in shape_text_lines(shape)]
    else:
        shape = next(
            (candidate for candidate in iter_shapes_recursive(slide.shapes) if str(candidate.shape_id) == str(node_id)),
            None,
        )
        if shape is None:
            fail(f"shape with nodeId {node_id!r} not found on slide {slide_number}")

        table_cell = anchor.get("tableCell")
        paragraph_index = anchor.get("paragraph")
        if table_cell is not None and paragraph_index is not None:
            fail("pptx anchor has both 'paragraph' and 'tableCell'; they address different things — pick one")
        if table_cell is not None:
            if not (getattr(shape, "has_table", False) and shape.has_table):
                fail(f"shape {node_id!r} is not a table but anchor has 'tableCell'")
            rows = list(shape.table.rows)
            row = require_index(table_cell.get("row"), "pptx anchor tableCell 'row'", 0)
            col = require_index(table_cell.get("col"), "pptx anchor tableCell 'col'", 0)
            if row >= len(rows) or col >= len(list(rows[row].cells)):
                fail(f"tableCell {table_cell!r} out of range for shape {node_id!r}")
            lines = [list(rows[row].cells)[col].text]
        elif paragraph_index is not None:
            if not shape.has_text_frame:
                fail(f"shape {node_id!r} has no text body but anchor has 'paragraph'")
            paragraphs = shape.text_frame.paragraphs
            # Bounded on both sides: a negative ordinal would index from the end and quietly return a
            # paragraph nobody asked for.
            paragraph_index = require_index(paragraph_index, "pptx anchor 'paragraph'", 0)
            if paragraph_index >= len(paragraphs):
                fail(f"paragraph {paragraph_index} out of range (shape has {len(paragraphs)} paragraphs)")
            lines = [paragraphs[paragraph_index].text]
        else:
            lines = shape_text_lines(shape)

    if out_format in ("txt", "md"):
        out_path.write_text("\n".join(lines) + "\n", encoding="utf-8")
    else:
        fail(f"unsupported output format for pptx source: {out_format!r} (use txt or md)")


EXTRACTORS = {"xlsx": extract_xlsx, "docx": extract_docx, "pdf": extract_pdf, "pptx": extract_pptx}


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("--file", required=True, help="absolute path of the source document (read-only)")
    parser.add_argument("--anchor", required=True, help="anchor JSON, e.g. '{\"format\":\"xlsx\",...}'")
    parser.add_argument("--out", required=True, help="absolute path of the NEW file to create; must not exist")
    args = parser.parse_args()

    src = Path(args.file)
    out_path = Path(args.out)
    # Both paths are documented, and schema-validated upstream, as absolute. Accepting a relative one
    # silently resolves it against whatever working directory the agent happens to be in.
    for label, candidate in (("--file", src), ("--out", out_path)):
        if not candidate.is_absolute():
            fail(f"{label} must be an absolute path: {str(candidate)!r}")
    if not src.is_file():
        fail(f"source file not found: {src}")
    if out_path.resolve() == src.resolve():
        fail("output path must differ from the source file — the source is never modified")
    if out_path.exists():
        fail(f"output path already exists: {out_path} — pick a fresh name instead of overwriting")

    try:
        anchor = json.loads(args.anchor)
    except json.JSONDecodeError as error:
        fail(f"anchor is not valid JSON: {error}")
    anchor_format = anchor.get("format")
    extractor = EXTRACTORS.get(anchor_format)
    if extractor is None:
        fail(f"unsupported anchor format: {anchor_format!r} (use xlsx, docx, pdf, or pptx)")

    out_format = out_path.suffix.lstrip(".").lower()
    if not out_format:
        fail("output path needs an extension so the output format can be inferred")

    if anchor_format in ("xlsx", "docx", "pptx"):
        preflight_zip(src)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    with atomic_output(out_path) as staging:
        extractor(src, anchor, staging, out_format)
    print(str(out_path))


if __name__ == "__main__":
    main()
